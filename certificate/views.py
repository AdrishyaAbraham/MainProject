from django.shortcuts import render, redirect, get_object_or_404
from django.http import HttpResponse
from django.core.mail import send_mail
from django.conf import settings
from django.contrib import messages
from django.contrib.auth.decorators import login_required
from .models import Event, Participant, Certificate_url
import pandas as pd
from .convter import ppt2pdf
from pptx import Presentation
from django.core.mail import send_mail, EmailMessage
import requests
import os
from django.views.decorators.cache import never_cache




@login_required
@never_cache
@login_required(login_url='login_page')
def create(request):
	if request.method == "POST":
		csv = request.FILES.get('csv')
		temp = request.FILES.get('template')
		event_name = request.POST.get('event_name')
		
		event = Event(user = request.user,
			event_name = event_name,
			csv_file = csv,
			template = temp)
		event.save()

		return redirect(f"/certificate/{event.id}/{event.slug}")

	return render(request, 'certificate/create_event.html')

# @login_required
# def delete_event(request, id, slug):
# 	event = Event.objects.filter(slug=slug, id=id).first()
# 	if event.user == request.user:
# 	    event.delete()
# 	return redirect('view_certificate_status')

@login_required
@never_cache
@login_required(login_url='login_page')
def track(request, id, slug):
	event = Event.objects.filter(slug=slug, id=id).first()
	if event.message:

		return render(request, 'certificate/track.html', {
			'event_name': event.event_name,
			'event_date': event.date,
			'participat_details': Participant.objects.filter(event=event)
			})

	prs = Presentation(event.template)
	st=""
	for slide in prs.slides:
		for shape in slide.shapes:
			if shape.has_text_frame:
				st = st + shape.text
				st = st + " "
	li = st.split()
	tags = []
	for i in li:
		if i[0] == "<" and i[-1] == ">":
			tags.append(i)
	print("Tags found in template:", tags)  # Debug message
	if request.method == "POST":
		email_col = request.POST.get('emails')
		subject = request.POST.get('subject')
		mess = request.POST.get('mess')
		values = [(tag, request.POST.get(f'type_{tag}'), request.POST.get(f'input_{tag}')) for tag in tags]	
		
		event.email_column = email_col
		event.message = mess
		event.subject = subject
		event.save()

		df=pd.read_csv(event.csv_file)

		df_len=df.shape
		i=0

		data = {
        	"client_id":"502998849458-0ebqeogcdvsr9lm7ug8jvdrenvbmfd6g.apps.googleusercontent.com",
        	"client_secret":"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",
        	"refresh_token": "1//04DIGqxYJAojMCgYIARAAGAQSNwF-L9Ir6kTYuxrU8N5P0dzblYRQSRCZxwVjR-zRJXukosCVJQ7zzCdgDbDqvAyQKaEZVwF1oqc",
        	'grant_type': 'refresh_token'
            }
		a = requests.post("https://www.googleapis.com/oauth2/v4/token", data)
		token = f"Bearer {dict(a.json()).get('access_token')}"
		li=["First","Second","Third"]
		while i < df_len[0]:
			prs = Presentation(event.template)
			j=""
			if i<9:
				j="00"
			elif i>=9 and i < 99 :
				j="0"
			
			for tag, v_type, value in values:
				for slide in prs.slides:
					for shape in slide.shapes:
						if shape.has_text_frame:
							if(shape.text.find(tag))!=-1:
									text_frame = shape.text_frame
									for paragraph in text_frame.paragraphs:
										for run in paragraph.runs:
											cur_text = run.text
											if v_type == 'text':
												new_text = cur_text.replace(tag, value)
											elif v_type == 'date':
												new_text = cur_text.replace(tag, '/'.join(value.split('-')[::-1]))
											elif v_type == 'csv':
												new_text = cur_text.replace(tag, df.loc[i,value])
											elif v_type == "auto":
												new_text = cur_text.replace(tag, value+"/"+j+str(i+1))
											else:
												pass
											run.text = new_text
											
			
			s_name = df.loc[i,event.email_column].split('@')[0]
			prs.save(s_name+".pptx")
			f_id = ppt2pdf(s_name+".pptx",s_name, token)
			r = requests.get(f"https://docs.google.com/presentation/d/{f_id}/export/pdf", allow_redirects=True)
			open(s_name+'.pdf', 'wb').write(r.content)

			try:
				mail = EmailMessage(subject,
					f"Hello, {s_name} \n{mess}",
					settings.EMAIL_HOST_USER,
					[df.loc[i,event.email_column]])
				mail.attach_file(s_name+'.pdf')
				mail.send()
				Participant(event=event, email=df.loc[i,event.email_column], status=True).save()
				os.remove(s_name+'.pdf')
				os.remove(s_name+".pptx")
				print(f"Email sent successfully to {df.loc[i, event.email_column]}")  # Debug message
			except Exception as e :
				Participant(event=event, email=df.loc[i,event.email_column], status=False).save()
				os.remove(s_name+'.pdf')
				os.remove(s_name+".pptx")
				print(f"Failed to send email to {df.loc[i, event.email_column]}")  # Debug message
                
		
			i=i+1

		messages.success(request, "Certificates Sent Successfuly !!")
		return redirect(f"/certificate/{event.id}/{event.slug}")


	return render(request, 'certificate/map_tags_of_template.html',{
		'columns': list(pd.read_csv(event.csv_file).columns),
		'tags': tags,
		})


@login_required
@never_cache
@login_required(login_url='login_page')
def view_certificate_status(request):
	return render(request, 'certificate/view_certificate_status.html',{
		'events': Event.objects.filter(user=request.user)
		})
@never_cache
@login_required(login_url='login_page')
def download_certificate(request, id, f_id):
    my_obj = get_object_or_404(Certificate_url,id = id, certificate_id = f_id)
    f_id = my_obj.certificate_id
    # return render(request, "certificate/download_certificate.html", {
    #    'url' : f"https://docs.google.com/presentation/d/{f_id}/export/pdf"
    #    })
    return HttpResponse(f"<a href = 'https://docs.google.com/presentation/d/{f_id}/export/pdf' target='_blank'>download</a>")
    return redirect(f"https://docs.google.com/presentation/d/{f_id}/export/pdf")

@login_required
@never_cache
@login_required(login_url='login_page')
def index(request):
	if request.method == "POST":
		csv = request.FILES.get('csv')
		temp = request.FILES.get('template')
		event_name = request.POST.get('event_name')
		event_type = request.POST.get('event_type')
		event_s_date = request.POST.get('s-date')
		event_e_date = request.POST.get('e-date')

		event = Event(user = request.user,
			event_name = event_name,
			event_type = event_type,
			starting_date = event_s_date,
			ending_date = event_e_date,
			csv_file = csv,
			template = temp)
		event.save()

		return redirect(f"/certificate/{event.id}/{event.slug}")

	return render(request, 'certificate/index.html', {'title': "Issue Certificates"})

# @login_required
# def delete_event(request, id, slug):
# 	event = Event.objects.filter(slug=slug, id=id).first()
# 	if event.user == request.user:
# 	    event.delete()
# 	return redirect('view_certificate_status')

# @login_required
# def track(request, id, slug):
# 	event = Event.objects.filter(slug=slug, id=id).first()
# 	if event.name_column:
# 		return render(request, 'certificate/track.html', {
# 			'event_name': event.event_name,
# 			'event_type': event.event_type,
# 			'event_date': event.starting_date,
# 			'participat_details': Participant.objects.filter(event=event),
# 			'title': f"Track Certificates of {event.event_name} {event.event_type}"
# 			})

# 	if request.method == "POST":
# 		name_col = request.POST.get('name_col')
# 		email_col = request.POST.get('email_col')
# 		org_col = request.POST.get('org_col')
# 		desi_col = request.POST.get('desi_col')
# 		sno_col = request.POST.get('sno_col')
# 		mess = request.POST.get('mess')

# 		event.name_column = name_col
# 		event.email_column = email_col
# 		event.org_column = org_col
# 		event.message = mess
# 		event.save()

# 		df = pd.read_csv(event.csv_file, encoding='latin1', errors='ignore')
# 		df_len = df.shape[0]
# 		i=0
# 		# Paste your google drive api details

# 		data = {
#         	"client_id":"YOUR-CLIENT-ID",
#         	"client_secret":"YOUR-CLIENT-SECRET",
#         	"refresh_token": "YOUR-REFRESH-TOKEN",
#         	'grant_type': 'refresh_token'
#             }
# 		a = requests.post("https://www.googleapis.com/oauth2/v4/token", data)
# 		token = f"Bearer {dict(a.json()).get('access_token')}"
# 		li=["First","Second","Third"]
# 		for i in range(df_len):
# 			s_name=df.loc[i,name_col]
# 			s_org=df.loc[i,org_col]
# 			s_email=df.loc[i,email_col]
# 			j=""
# 			if i<9:
# 				j="00"
# 			elif i>=9 and i < 99 :
# 				j="0"
# 			prs = Presentation(event.template)
# 			if desi_col != "Choose...":
# 				s_desi=df.loc[i,desi_col]
# 				for slide in prs.slides:
# 					for shape in slide.shapes:
# 						if shape.has_text_frame:
# 							if(shape.text.find("{{Full Name}}"))!=-1:
# 								text_frame = shape.text_frame
# 								for paragraph in text_frame.paragraphs:
# 									for run in paragraph.runs:
# 										cur_text = run.text
# 										new_text = cur_text.replace("{{Full Name}}", s_name+" , "+s_desi)
# 										run.text = new_text
# 			else:
# 				for slide in prs.slides:
# 					for shape in slide.shapes:
# 						if shape.has_text_frame:
# 							if(shape.text.find("{{Full Name}}"))!=-1:
# 								text_frame = shape.text_frame
# 								for paragraph in text_frame.paragraphs:
# 									for run in paragraph.runs:
# 										cur_text = run.text
# 										new_text = cur_text.replace("{{Full Name}}", s_name)
# 										run.text = new_text

# 			for slide in prs.slides:
# 				for shape in slide.shapes:
# 					if shape.has_text_frame:
# 						if(shape.text.find("{{certificate id}}"))!=-1:
# 								text_frame = shape.text_frame
# 								for paragraph in text_frame.paragraphs:
# 									for run in paragraph.runs:
# 										cur_text = run.text
# 										new_text = cur_text.replace("{{certificate id}}", "ISM/"+sno_col+"/"+j+str(i+1))
# 										run.text = new_text

# 			for slide in prs.slides:
# 				for shape in slide.shapes:
# 					if shape.has_text_frame:
# 						if(shape.text.find("{{organization}}"))!=-1:
# 								text_frame = shape.text_frame
# 								for paragraph in text_frame.paragraphs:
# 									for run in paragraph.runs:
# 										cur_text = run.text
# 										new_text = cur_text.replace("{{organization}}", s_org)
# 										run.text = new_text
# 			if event.event_type == "Quiz" and i < 3:
# 				for slide in prs.slides:
# 					for shape in slide.shapes:
# 						if shape.has_text_frame:
# 							if(shape.text.find("participated"))!=-1:
# 								text_frame = shape.text_frame
# 								for paragraph in text_frame.paragraphs:
# 									for run in paragraph.runs:
# 										cur_text = run.text
# 										new_text = cur_text.replace("participated", "secured "+str(li[i])+" position")
# 										run.text = new_text
# 			for slide in prs.slides:
# 				for shape in slide.shapes:
# 					if shape.has_text_frame:
# 						if(shape.text.find("{{Event Type}}"))!=-1:
# 								text_frame = shape.text_frame
# 								for paragraph in text_frame.paragraphs:
# 									for run in paragraph.runs:
# 										cur_text = run.text
# 										new_text = cur_text.replace("{{Event Type}}", event.event_type)
# 										run.text = new_text

# 			for slide in prs.slides:
# 				for shape in slide.shapes:
# 					if shape.has_text_frame:
# 						if(shape.text.find("{{topic}}"))!=-1:
# 								text_frame = shape.text_frame
# 								for paragraph in text_frame.paragraphs:
# 									for run in paragraph.runs:
# 										cur_text = run.text
# 										new_text = cur_text.replace("{{topic}}", event.event_name)
# 										run.text = new_text
# 			for slide in prs.slides:
# 				for shape in slide.shapes:
# 					if shape.has_text_frame:
# 						if(shape.text.find("{{Date}}"))!=-1:
# 								text_frame = shape.text_frame
# 								for paragraph in text_frame.paragraphs:
# 									for run in paragraph.runs:
# 										cur_text = run.text
# 										if event.starting_date == event.ending_date:
# 											new_text = cur_text.replace("{{Date}}", str(event.starting_date))
# 										else:
# 											new_text = cur_text.replace("{{Date}}", str(event.starting_date))
# 										run.text = new_text
# 			#prs.save(".\certificate\certificate.pptx")
# 			prs.save("certificate.pptx")
# 			f_id = ppt2pdf("certificate.pptx",s_name, token)
# 			r = requests.get(f"https://docs.google.com/presentation/d/{f_id}/export/pdf", allow_redirects=True)
# 			open(s_name+'.pdf', 'wb').write(r.content)

# 			try:
# 				mail = EmailMessage(f"Certificate for {event.event_name} {event.event_type}",
# 					f"Hello, {s_name} \n{mess}",
# 					settings.EMAIL_HOST_USER,
# 					[s_email])
# 				mail.attach_file(s_name+'.pdf')
# 				mail.send()
# 				Participant(event=event, name=s_name, email=s_email, org=s_org, status=True).save()
# 				os.remove(s_name+'.pdf')
# 				os.remove("certificate.pptx")
# 			except:
# 				Participant(event=event, name=s_name, email=s_email, org=s_org, status=False).save()

# 			i=i+1

# 		messages.success(request, "Certificates Sent Successfuly !!")
# 		return redirect(f"/certificate/{event.id}/{event.slug}")

# 	return render(request, 'certificate/from2.html',{
# 		'columns': list(pd.read_csv(event.csv_file).columns),
# 		'title': 'Send Certificates'
# 		})


# @login_required
# def view_certificate_status(request):
# 	return render(request, 'certificate/view_certificate_status.html',{
# 		'events': Event.objects.filter(user=request.user),
# 		'title': "All Events"
# 		})

@never_cache
@login_required(login_url='login_page')
@login_required
def upload_template(request):
	return render(request, 'certificate/upload_template.html')

